
from abc import ABC, abstractmethod
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import ttkbootstrap as tb
from PIL import Image
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from pptx import Presentation
import pytesseract
import io
import os
import fitz
from pdf2image import convert_from_path
from docx import Document
import comtypes.client
import mysql.connector
import pdf2docx
from pdf2docx import Converter

# Abstract Base Class for Abstraction
class AbstractPDFConverterApp(ABC):
    @abstractmethod
    def pdf_to_ppt(self):
        pass

    @abstractmethod
    def pdf_to_word(self):
        pass

    @abstractmethod
    def jpg_to_pdf(self):
        pass

    @abstractmethod
    def word_to_pdf(self):
        pass


class PDFConverterApp(AbstractPDFConverterApp):
    def __init__(self, root):
        self.__root = root  # Encapsulating root as a private variable
        self.__root.title("PDF Converter")
        self.__root.geometry("800x600")
        self.__style = tb.Style(theme="cosmo")  # Choose a visually appealing theme

        # Set background color
        self.__root.configure(bg="#2C3E50")

        # Database connection setup
        self.__db_connection()

        # Welcome Page
        self.show_welcome_page()

# Encapsulated Database Connection
    def __db_connection(self):
        try:
            self.__conn = mysql.connector.connect(
                host="localhost",
                user="root",
                password="Sabbir@190921",
                database="pdf_converter_db",  # Ensure this database exists
            )
            self.__cursor = self.__conn.cursor()
        except mysql.connector.Error as e:
            messagebox.showerror(
                "Database Error", f"Failed to connect to database: {e}"
            )

 # Getter for Database Cursor
    @property
    def cursor(self):
        return self.__cursor

# Welcome Page
    def show_welcome_page(self):
        for widget in self.__root.winfo_children():
            widget.destroy()

        welcome_label = ttk.Label(
            self.__root,
            text="Welcome to PDF Converter",
            font=("Helvetica", 24, "bold"),
            foreground="white",
            background="#2C3E50",
        )
        welcome_label.pack(pady=20)

        sign_in_button = ttk.Button(
            self.__root,
            text="Sign In",
            command=self.show_sign_in_page,
            width=20,
            bootstyle="primary",
        )
        sign_in_button.pack(pady=10)

        sign_up_button = ttk.Button(
            self.__root,
            text="Sign Up",
            command=self.show_sign_up_page,
            width=20,
            bootstyle="success",
        )
        sign_up_button.pack(pady=10)
#sign in page
    def show_sign_in_page(self):
        for widget in self.__root.winfo_children():
            widget.destroy()

        ttk.Label(
            self.__root,
            text="Sign In",
            font=("Helvetica", 20, "bold"),
            foreground="green",
        ).pack(pady=20)

        username_label = ttk.Label(self.__root, text="Username:")
        username_label.pack(pady=5)
        username_entry = ttk.Entry(self.__root)
        username_entry.pack(pady=5)

        password_label = ttk.Label(self.__root, text="Password:")
        password_label.pack(pady=5)
        password_entry = ttk.Entry(self.__root, show="*")
        password_entry.pack(pady=5)

        sign_in_button = ttk.Button(
            self.__root,
            text="Sign In",
            command=lambda: self.validate_user(username_entry, password_entry),
            width=20,
            bootstyle="primary",
        )
        sign_in_button.pack(pady=20)

        back_button = ttk.Button(
            self.__root, text="Back", command=self.show_welcome_page, bootstyle="danger"
        )
        back_button.pack()
#valid user
    def validate_user(self, username_entry, password_entry):
        username = username_entry.get()
        password = password_entry.get()
        self.cursor.execute(
            "SELECT * FROM users WHERE username = %s AND password = %s",
            (username, password),
        )
        if self.cursor.fetchone():
            messagebox.showinfo("Success", "Login successful!")
            self.show_feature_page()
        else:
            messagebox.showerror("Error", "Invalid username or password!")
#sign up page
    def show_sign_up_page(self):
        for widget in self.__root.winfo_children():
            widget.destroy()

        ttk.Label(
            self.__root,
            text="Sign Up",
            font=("Helvetica", 20, "bold"),
            foreground="green",
        ).pack(pady=20)

        username_label = ttk.Label(self.__root, text="Username:")
        username_label.pack(pady=5)
        username_entry = ttk.Entry(self.__root)
        username_entry.pack(pady=5)

        password_label = ttk.Label(self.__root, text="Password:")
        password_label.pack(pady=5)
        password_entry = ttk.Entry(self.__root, show="*")
        password_entry.pack(pady=5)

        sign_up_button = ttk.Button(
            self.__root,
            text="Sign Up",
            command=lambda: self.register_user(username_entry, password_entry),
            width=20,
            bootstyle="success",
        )
        sign_up_button.pack(pady=20)

        back_button = ttk.Button(
            self.__root, text="Back", command=self.show_welcome_page, bootstyle="danger"
        )
        back_button.pack()
#register user
    def register_user(self, username_entry, password_entry):
        username = username_entry.get()
        password = password_entry.get()
        if not username or not password:
            messagebox.showerror("Input Error", "All fields are required!")
            return
        try:
            self.cursor.execute(
                "INSERT INTO users (username, password) VALUES (%s, %s)",
                (username, password),
            )
            self.__conn.commit()
            messagebox.showinfo("Success", "User registered successfully!")
            self.show_sign_in_page()
        except mysql.connector.Error as e:
            messagebox.showerror("Database Error", f"Failed to register user: {e}")
#Feature page
    def show_feature_page(self):
        for widget in self.__root.winfo_children():
            widget.destroy()

        ttk.Label(
            self.__root,
            text="Choose a Feature",
            font=("Times New Roman", 20, "bold"),
            foreground="white",
            background="#2C3E50",
        ).pack(pady=20)

        features_left = [
            ("PDF to PPT", self.pdf_to_ppt),
            ("PDF to Word", self.pdf_to_word)
        ]

        features_right = [
            ("JPG to PDF", self.jpg_to_pdf),
            ("Word to PDF", self.word_to_pdf)
        ]

        left_frame = ttk.Frame(self.__root)
        left_frame.pack(side="left", padx=50)

        right_frame = ttk.Frame(self.__root)
        right_frame.pack(side="right", padx=50)

        for feature_name, feature_command in features_left:
            ttk.Button(
                left_frame,
                text=feature_name,
                command=feature_command,
                width=30,
                bootstyle="info",
            ).pack(pady=10)

        for feature_name, feature_command in features_right:
            ttk.Button(
                right_frame,
                text=feature_name,
                command=feature_command,
                width=30,
                bootstyle="info",
            ).pack(pady=10)
#Pdt to ppt
    # Implementation of Abstract Methods
    def pdf_to_ppt(self):
        # Open file dialog to select PDF
        file_path = filedialog.askopenfilename(
            filetypes=[("PDF Files", "*.pdf")], title="Select PDF File"
        )

        if not file_path:
            return  # If no file is selected, return

        # Create a presentation object
        ppt = Presentation()

        # Open the PDF using PyMuPDF (fitz)
        pdf_document = fitz.open(file_path)

        # Loop through each page of the PDF and convert it to an image
        for page_num in range(pdf_document.page_count):
            # Extract each page as an image (pixmap)
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap()

            # Convert pixmap to a PIL Image object
            img_data = pix.tobytes("png")  # You can change format to 'jpeg' if needed
            img = Image.open(io.BytesIO(img_data))

            # Save image to a temporary file (for presentation)
            img_path = f"page_{page_num + 1}.png"
            img.save(img_path)

            # Create a new slide in the presentation
            slide_layout = ppt.slide_layouts[5]  # Blank slide layout
            slide = ppt.slides.add_slide(slide_layout)

            # Add the image to the slide
            slide.shapes.add_picture(
                img_path, 0, 0, width=ppt.slide_width, height=ppt.slide_height
            )

            # Optionally, delete the temporary image file
            img.close()

        # Save the presentation
        output_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint Files", "*.pptx")],
            title="Save PowerPoint File",
        )

        if output_path:
            ppt.save(output_path)
            messagebox.showinfo("Success", "PDF successfully converted to PowerPoint!")



    def pdf_to_word(self):
        file_path = filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")])
        if not file_path:
            return
        output_path = filedialog.asksaveasfilename(
            defaultextension=".docx", filetypes=[("Word Documents", "*.docx")]
        )
        if not output_path:
            return
        try:
            converter = Converter(file_path)
            # Convert the PDF to a Word document and save it to the given file path
            converter.convert(
                output_path, start=0, end=None
            )  # You can specify page range if needed
            converter.close()
            print(f"Conversion complete! Word document saved as: {output_path}")
        except Exception as e:
            messagebox.showerror(
                "Error", f"An error occurred while converting PDF to Word: {e}"
            )


    def jpg_to_pdf(self):
         # Ask for the JPG file location
        file_path = filedialog.askopenfilename(filetypes=[("JPG Files", "*.jpg")])

        if not file_path:
            return

        # Ask for output PDF file location
        output_path = filedialog.asksaveasfilename(
            defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")]
        )

        if not output_path:
            return

        try:
            img = Image.open(file_path)
            pdf_path = output_path
            img.convert("RGB").save(pdf_path, "PDF")
            messagebox.showinfo("Success", "JPG successfully converted to PDF!")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred during conversion: {e}")

    
    def word_to_pdf(self):
        """
        Converts a Word document (.docx) to a PDF file.
        """
        # Ask for Word file location
        file_path = filedialog.askopenfilename(
            title="Select Word Document",
            filetypes=[("Word Documents", "*.docx")]
        )

        if not file_path:
            messagebox.showwarning("Warning", "No Word document selected!")
            return

        # Ask for output PDF file location
        output_path = filedialog.asksaveasfilename(
            title="Save as PDF",
            defaultextension=".pdf",
            filetypes=[("PDF Files", "*.pdf")]
        )

        if not output_path:
            messagebox.showwarning("Warning", "No output file path provided!")
            return

        try:
            # Verify file and folder paths
            if not os.path.exists(file_path):
                raise FileNotFoundError("The selected Word document does not exist.")

            output_dir = os.path.dirname(output_path)
            if not os.path.isdir(output_dir):
                raise FileNotFoundError("The specified output folder does not exist.")

            # Initialize Word COM object
            word = comtypes.client.CreateObject("Word.Application")
            word.Visible = 0  # Set to 1 for debugging to make Word visible

            # Open the Word document
            doc = word.Documents.Open(file_path)

            # Save the document as a PDF
            doc.SaveAs(output_path, FileFormat=17)  # 17 is the format for PDF
            
            # Close the document and Word application
            doc.Close()
            word.Quit()

            # Success message
            messagebox.showinfo(
                "Success",
                f"Word document successfully converted to PDF!\nSaved at: {output_path}"
            )

        except Exception as e:
            # Error handling
            messagebox.showerror("Error", f"An error occurred during conversion: {e}")

        finally:
            try:
                word.Quit()  # Ensure Word quits even on error
            except:
                pass

if __name__ == "__main__":
    root = tb.Window(themename="cosmo")
    app = PDFConverterApp(root)
    root.mainloop()
